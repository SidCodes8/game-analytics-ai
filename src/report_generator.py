from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import seaborn as sns
import io
import base64
import pandas as pd
from typing import Dict, List
import os

class ReportGenerator:
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.darkblue
        )
        self.heading_style = ParagraphStyle(
            'CustomHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            textColor=colors.darkblue
        )
    
    def generate_pdf_report(self, analytics_data: Dict, insights: str, 
                           output_path: str = "reports/analytics_report.pdf") -> str:
        """Generate comprehensive PDF report"""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        story = []
        
        # Title
        story.append(Paragraph("Gaming Analytics Report", self.title_style))
        story.append(Spacer(1, 12))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", self.heading_style))
        story.append(Paragraph("This report provides comprehensive insights into player behavior, revenue trends, and actionable recommendations for improving game monetization and retention.", self.styles['Normal']))
        story.append(Spacer(1, 12))
        
        # Key Metrics
        story.append(Paragraph("Key Metrics", self.heading_style))
        
        if 'total_revenue' in analytics_data:
            story.append(Paragraph(f"<b>Total Revenue:</b> ${analytics_data['total_revenue']:,.2f}", self.styles['Normal']))
        
        if 'avg_arppu' in analytics_data:
            story.append(Paragraph(f"<b>Average ARPPU:</b> ${analytics_data['avg_arppu']:.2f}", self.styles['Normal']))
        
        if 'avg_arpdau' in analytics_data:
            story.append(Paragraph(f"<b>Average ARPDAU:</b> ${analytics_data['avg_arpdau']:.2f}", self.styles['Normal']))
        
        story.append(Spacer(1, 12))
        
        # User Segments
        if 'segment_summary' in analytics_data:
            story.append(Paragraph("User Segments", self.heading_style))
            segment_data = analytics_data['segment_summary']
            
            # Convert segment summary to table
            table_data = [['Segment', 'Users', 'Total Spend', 'Avg Spend']]
            for segment in segment_data.index:
                row = [
                    segment.title(),
                    str(segment_data.loc[segment, ('user_id', 'count')] if ('user_id', 'count') in segment_data.columns else 'N/A'),
                    f"${segment_data.loc[segment, ('total_spend', 'sum')]:,.2f}" if ('total_spend', 'sum') in segment_data.columns else 'N/A',
                    f"${segment_data.loc[segment, ('total_spend', 'mean')]:,.2f}" if ('total_spend', 'mean') in segment_data.columns else 'N/A'
                ]
                table_data.append(row)
            
            table = Table(table_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 14),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            story.append(table)
            story.append(Spacer(1, 12))
        
        # AI Insights
        story.append(Paragraph("AI-Generated Insights & Recommendations", self.heading_style))
        story.append(Paragraph(insights, self.styles['Normal']))
        story.append(Spacer(1, 12))
        
        # Anomalies
        if 'anomalies' in analytics_data and analytics_data['anomalies']:
            story.append(Paragraph("Detected Anomalies", self.heading_style))
            for anomaly in analytics_data['anomalies'][:5]:  # Show top 5 anomalies
                story.append(Paragraph(
                    f"<b>{anomaly['date']}</b>: {anomaly['type'].title()} in {anomaly['metric']} "
                    f"(Value: {anomaly['value']:.2f}, Expected: {anomaly['expected']:.2f})",
                    self.styles['Normal']
                ))
            story.append(Spacer(1, 12))
        
        # Build PDF
        doc.build(story)
        return output_path
    
    def generate_pptx_report(self, analytics_data: Dict, insights: str,
                           output_path: str = "reports/analytics_presentation.pptx") -> str:
        """Generate PowerPoint presentation"""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Gaming Analytics Report"
        subtitle.text = "Data-Driven Insights for Player Engagement & Monetization"
        
        # Key Metrics slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "Key Performance Metrics"
        
        tf = body.text_frame
        if 'total_revenue' in analytics_data:
            p = tf.add_paragraph()
            p.text = f"Total Revenue: ${analytics_data['total_revenue']:,.2f}"
            p.level = 0
        
        if 'avg_arppu' in analytics_data:
            p = tf.add_paragraph()
            p.text = f"Average ARPPU: ${analytics_data['avg_arppu']:.2f}"
            p.level = 0
        
        if 'avg_arpdau' in analytics_data:
            p = tf.add_paragraph()
            p.text = f"Average ARPDAU: ${analytics_data['avg_arpdau']:.2f}"
            p.level = 0
        
        # AI Insights slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "AI-Generated Insights"
        
        tf = body.text_frame
        tf.text = insights[:500] + "..." if len(insights) > 500 else insights
        
        # Recommendations slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "Recommendations"
        
        tf = body.text_frame
        recommendations = [
            "Implement personalized offers for high-value segments",
            "Focus retention efforts on at-risk players",
            "Optimize pricing based on user behavior patterns",
            "Launch targeted campaigns for inactive users"
        ]
        
        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.level = 0
        
        prs.save(output_path)
        return output_path
    
    def create_chart_image(self, data: pd.DataFrame, chart_type: str, 
                          title: str, x_col: str, y_col: str) -> str:
        """Create chart image for reports"""
        plt.figure(figsize=(10, 6))
        
        if chart_type == 'line':
            plt.plot(data[x_col], data[y_col])
        elif chart_type == 'bar':
            plt.bar(data[x_col], data[y_col])
        elif chart_type == 'scatter':
            plt.scatter(data[x_col], data[y_col])
        
        plt.title(title)
        plt.xlabel(x_col.replace('_', ' ').title())
        plt.ylabel(y_col.replace('_', ' ').title())
        plt.xticks(rotation=45)
        plt.tight_layout()
        
        # Save to bytes
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
        img_buffer.seek(0)
        
        # Convert to base64 for embedding
        img_str = base64.b64encode(img_buffer.read()).decode()
        plt.close()
        
        return img_str